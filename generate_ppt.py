import os
import re
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from pptx.enum.text import MSO_AUTO_SIZE
from apis.groq_api import GroqAPIClient
from utils import load_config
from crawlers.serpapi_image import SerpAPIImageFetcher, UnsplashFallback
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR


def load_theme():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    theme_path = os.path.join(base_dir, "theme0.pptx")
    return Presentation(theme_path)


def ensure_pptx_compatible_image(img_path):
    try:
        with Image.open(img_path) as img:
            img.verify()

        with Image.open(img_path) as img:
            if img.format == "WEBP":
                new_path = img_path.rsplit(".", 1)[0] + ".png"
                img.convert("RGB").save(new_path, "PNG")
                return new_path

        return img_path

    except Exception:
        return None


# ----------------------------------------------------
# Slide Creator Functions
# ----------------------------------------------------


from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


def create_title_slide(ppt, title, subtitle):
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])

    # ---------------- TITLE ----------------
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.9), Inches(9.2), Inches(1.8)
    )

    tf = title_box.text_frame
    tf.text = title
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.auto_size = MSO_AUTO_SIZE.NONE

    p = tf.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.bold = True

    # base font
    p.font.size = Pt(46)

    # Auto shrink for long titles
    if len(title) > 35:
        p.font.size = Pt(42)
    if len(title) > 55:
        p.font.size = Pt(38)
    if len(title) > 70:
        p.font.size = Pt(34)

    # separator line
    line = slide.shapes.add_shape(
        autoshape_type_id=1,
        left=Inches(1.2),
        top=Inches(2.25),
        width=Inches(7.6),
        height=Pt(0.5)
    )
    line.line.width = Pt(1.5)

    # ---------------- SUBTITLE ----------------
    if subtitle:

        text = subtitle.replace("•", " ").strip()
        max_len = 310

        text = subtitle.replace("•", " ").strip()

        if len(text) > max_len:
            trimmed = text[:max_len]

            # try to end at full stop
            last_period = trimmed.rfind(".")
            last_space = trimmed.rfind(" ")

            if last_period != -1 and last_period > 120:
                trimmed = trimmed[:last_period]
            else:
                trimmed = trimmed[:last_space]

            text = trimmed + "..."

        subtitle_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.9), Inches(8.2), Inches(3.0)
        )

        st = subtitle_box.text_frame
        st.word_wrap = True
        st.text = text

        sp = st.paragraphs[0]
        sp.font.size = Pt(19)
        sp.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER




def create_content_slide(ppt, title, content, is_first=False):
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(9), Inches(1.1 if is_first else 0.9)
    )
    tf_title = title_box.text_frame
    tf_title.auto_size = MSO_AUTO_SIZE.NONE
    tf_title.text = title
    tf_title.paragraphs[0].font.size = Pt(28 if is_first else 26)
    tf_title.paragraphs[0].font.bold = True

    content_box = slide.shapes.add_textbox(
        Inches(0.7),
        Inches(1.6),
        Inches(8.6),
        Inches(5.2 if is_first else 4.8)
    )

    tf = content_box.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True
    tf.clear()

    for line in [l.strip("• ").strip() for l in content.split("\n") if l.strip()]:
        p = tf.add_paragraph()
        p.text = "• " + line
        p.font.size = Pt(21 if is_first else 20)

        if is_first:
            p.line_spacing = 1.35
            p.space_after = Pt(10)
        else:
            p.line_spacing = 1.15
            p.space_after = Pt(6)


def create_image_slide(ppt, title, content, img_path):
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])

    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.3), Inches(9), Inches(0.8)
    )
    t = title_box.text_frame
    t.text = title

    try:
        t.paragraphs[0].font.size = Pt(28)
        t.paragraphs[0].font.bold = True
    except:
        pass

    content_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(1.2), Inches(4.2), Inches(4.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    bullets = [line.strip("• ").strip() for line in content.split("\n") if line.strip()]

    if hasattr(tf, "clear"):
        tf.clear()

    small_font = Pt(16)

    for line in bullets:
        p = tf.add_paragraph()
        p.text = "• " + line
        try:
            p.font.size = small_font
        except:
            pass

    try:
        img = Image.open(img_path)
        img_w, img_h = img.size

        max_w = Inches(4)
        max_h = Inches(4)
        aspect = img_w / img_h

        if aspect >= 1:
            width = max_w
            height = width / aspect
        else:
            height = max_h
            width = height * aspect

        slide.shapes.add_picture(
            img_path,
            Inches(5.4),
            Inches(1.2),
            width=width,
            height=height
        )

    except Exception as e:
        print("Image failed:", e)


# ----------------------------------------------------
# Parse AI text
# ----------------------------------------------------
def parse_slides(text):
    chunks = text.split("[SLIDEBREAK]")
    final_slides = []
    seen_titles = set()

    for s in chunks:
        title_match = re.search(r"\[TITLE\](.*?)\[/TITLE\]", s, re.S)
        content_match = re.search(r"\[CONTENT\](.*?)\[/CONTENT\]", s, re.S)

        title = title_match.group(1).strip() if title_match else ""
        content = content_match.group(1).strip() if content_match else ""

        if not title and not content:
            continue

        norm = title.lower().strip()
        if norm in seen_titles:
            continue

        seen_titles.add(norm)
        final_slides.append((title, content))

    return final_slides


# ----------------------------------------------------
# MAIN GENERATION FUNCTION (UNCHANGED)
# ----------------------------------------------------
def generate_ppt(topic, api_name, model_name, num_slides, api_key,serp_api_key=None):
    config = load_config()
    save_dir = config.get("save_location", "generated_ppts")
    os.makedirs(save_dir, exist_ok=True)

    ppt = load_theme()
    serp = SerpAPIImageFetcher()

    # if user typed a key — override config key
    if serp_api_key:
        serp.api_key = serp_api_key

    unsplash = UnsplashFallback()

    client = GroqAPIClient(api_key, model_name)

    prompt = f"""
    Create a professional PowerPoint outline with EXACTLY {num_slides} slides
    about "{topic}", using this structure:

    [TITLE]Slide Title[/TITLE]
    [CONTENT]
    • 3–5 expert-level bullet points.
    • Business-professional explanation.
    [/CONTENT]
    [SLIDEBREAK]
    """

    raw = client.generate(prompt)
    slides = parse_slides(raw)

    if slides and slides[0][0].strip() == "":
        slides[0] = (topic, slides[0][1])

    slides.append(("Thank You", ""))

    mid_indices = list(range(1, len(slides) - 1))
    slides_with_images = set(random.sample(mid_indices, max(1, len(mid_indices) // 2)))

    for i, (title, content) in enumerate(slides):

        if i == 0:
            lines = [
                l.replace("•", "").strip()
                for l in content.split("\n")
                if l.strip()
                ]

            # Use ALL text — let PowerPoint wrap naturally
            paragraph = " ".join(lines)

            create_title_slide(ppt, title, paragraph)
            continue

        if i == len(slides) - 1:
            create_content_slide(ppt, title, content)
            continue

        if i in slides_with_images:
            query = re.sub(r"[^a-zA-Z0-9 ]", "", f"{topic} {title}")
            img_path = serp.get_image(query, save_dir) or unsplash.get_image(query, save_dir)

            if img_path and os.path.exists(img_path):
                img_path = ensure_pptx_compatible_image(img_path)

                if img_path:
                    try:
                        create_image_slide(ppt, title, content, img_path)
                        continue
                    except Exception as e:
                        print("Image insertion failed:", e)

        create_content_slide(ppt, title, content)

    clean_title = re.sub(r"[^a-zA-Z0-9 ]", "", topic).replace(" ", "_")
    path = os.path.join(save_dir, f"{clean_title}.pptx")
    ppt.save(path)

    return f"PPT Created Successfully!\nSaved at: {path}"

# ----------------------------------------------------
# API WRAPPER (RETURN ONLY FILE PATH)
# ----------------------------------------------------
def generate_ppt_api(topic, api_name, model_name, num_slides, api_key, serp_api_key=None):
    """
    Wrapper for API usage — keeps original logic,
    but returns ONLY the ppt file path.
    """

    result = generate_ppt(
        topic=topic,
        api_name=api_name,
        model_name=model_name,
        num_slides=num_slides,
        api_key=api_key,
        serp_api_key=serp_api_key
    )

    for line in result.splitlines():
        if line.lower().startswith("saved at"):
            return line.split(":", 1)[1].strip()

    raise ValueError("Could not extract PPT path.")
